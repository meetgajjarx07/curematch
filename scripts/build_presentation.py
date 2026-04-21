#!/usr/bin/env python3
"""
Generate CureMatch_Presentation.pptx — editorial Apple-style deck
with embedded data visualizations and full speaker notes for every slide.

18 slides · 12-minute talk · three presenters:

   1. Cover
   2. The problem                            [narrator]
   3. The corpus at a glance                 [narrator]              ← chart
   4. What we built                          [Aashish opens]
   5. Architecture · LLMs read · Rules · Humans
   6. Team · three engineers · three LLMs
   7. Aashish · The Match Explainer
   8. Aashish · Prompt engineering
   9. Meet · The whole application
  10. Meet · The Chat Agent
  11. Meet · LoRA fine-tune                   ← loss curve
  12. Daksh · The LLM Parser
  13. Daksh · Parser coverage                 ← chart
  14. Daksh · Scoring engine
  15. End-to-end funnel                       ← chart
  16. In numbers
  17. Limitations
  18. Closing

Each slide has speaker notes — open Presenter View in PowerPoint
(play ▸ Presenter View) to read them while presenting.

Run:
    python3 scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PROJECT_ROOT = Path(__file__).resolve().parent.parent
CHARTS_DIR = PROJECT_ROOT / "data" / "charts"


# ─── Design tokens ──────────────────────────────────────────────
INK       = RGBColor(0x1D, 0x1D, 0x1F)
INK_SOFT  = RGBColor(0x3B, 0x3B, 0x41)
INK_MUTE  = RGBColor(0x6E, 0x6E, 0x73)
INK_FAINT = RGBColor(0x86, 0x86, 0x8B)
LINE      = RGBColor(0xD2, 0xD2, 0xD7)
LINE_SOFT = RGBColor(0xE8, 0xE8, 0xED)
PAPER     = RGBColor(0xFB, 0xFB, 0xFD)
PAPER_ALT = RGBColor(0xF5, 0xF5, 0xF7)
DEEP      = RGBColor(0x05, 0x07, 0x0F)
ACCENT    = RGBColor(0x00, 0x71, 0xE3)
ACCENT_LT = RGBColor(0x29, 0x97, 0xFF)
SUCCESS   = RGBColor(0x30, 0xD1, 0x58)
WARNING   = RGBColor(0xFF, 0x9F, 0x0A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT_SANS = "Helvetica Neue"
FONT_MONO = "JetBrains Mono"

SLIDE_W_IN = 10.667
SLIDE_H_IN = 7.5
TOTAL_SLIDES = 19


# ─── Primitives ─────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text(slide, left, top, width, height, body,
         *, size=16, color=INK, bold=False, italic=False,
         font=FONT_SANS, align="left", anchor="top",
         line_spacing=1.2, letter_spacing=None):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        if letter_spacing is not None:
            r._r.get_or_add_rPr().set("spc", str(letter_spacing))
    return tb


def rect(slide, left, top, width, height, fill,
         *, line=None, rounded=False, radius=0.06):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(
        shp, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    if rounded:
        s.adjustments[0] = radius
    return s


def image(slide, path: Path, *, left, top, width=None, height=None):
    if not path.exists():
        w = width or 5
        h = height or 3
        rect(slide, left, top, w, h, PAPER_ALT, line=LINE_SOFT, rounded=True)
        text(slide, left, top + h / 2 - 0.2, w, 0.4,
             f"[ {path.name} not found — run the generator ]",
             size=10, color=INK_FAINT, align="center", italic=True)
        return None
    kw = {}
    if width is not None:
        kw["width"] = Inches(width)
    if height is not None:
        kw["height"] = Inches(height)
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kw)


# ─── Slide chrome ───────────────────────────────────────────────
def eyebrow(slide, label: str, *, top=0.55, color=INK_MUTE):
    text(slide, 0.6, top, 9.5, 0.3, label.upper(),
         size=10, color=color, bold=True, letter_spacing=250)


def footer(slide, page: int, *, dark=False):
    c = RGBColor(0x86, 0x86, 0x8B) if dark else INK_FAINT
    text(slide, 0.5, 7.08, 3, 0.3, "CureMatch",
         size=10, color=c, letter_spacing=100)
    text(slide, 9.5, 7.08, 0.9, 0.3, f"{page:02d} / {TOTAL_SLIDES:02d}",
         size=10, color=c, align="right", font=FONT_MONO)


def title(slide, body: str, *, top=1.1, size=46, color=INK, height=1.6,
          bold=True, line_spacing=1.05, letter_spacing=-28, align="left"):
    text(slide, 0.6, top, SLIDE_W_IN - 1.2, height, body,
         size=size, color=color, bold=bold, letter_spacing=letter_spacing,
         line_spacing=line_spacing, align=align)


def notes(slide, body: str):
    """Attach speaker notes (presenter-only text) to this slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    for i, para in enumerate(body.strip().split("\n")):
        if i == 0:
            p = notes_tf.paragraphs[0]
        else:
            p = notes_tf.add_paragraph()
        p.text = para


def talkpoints(slide, left, top, width, height, bullets):
    """A muted, readable list of bullets the presenter can glance at."""
    rect(slide, left, top, width, height, PAPER_ALT,
         line=LINE_SOFT, rounded=True)
    text(slide, left + 0.25, top + 0.18, width - 0.5, 0.3,
         "TALKING POINTS", size=9, color=INK_FAINT,
         bold=True, letter_spacing=200)
    body = "\n".join(f"•  {b}" for b in bullets)
    text(slide, left + 0.25, top + 0.5, width - 0.5, height - 0.55, body,
         size=11, color=INK_SOFT, line_spacing=1.5)


# ─── Build ──────────────────────────────────────────────────────
def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════════
    # 01 · COVER
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)

    text(s, 0.6, 0.55, 10, 0.3, "FINAL GROUP PRESENTATION · APRIL 2026",
         size=10, color=INK_MUTE, bold=True, letter_spacing=300)

    text(s, 0.6, 2.0, 10, 1.5, "CureMatch.",
         size=128, color=INK, bold=True, letter_spacing=-60)

    text(s, 0.6, 3.8, 10, 0.9,
         "Matching clinical trials with large language models.",
         size=26, color=INK_MUTE)

    rect(s, 0.6, 5.4, 0.5, 0.04, ACCENT)
    text(s, 0.6, 5.65, 10, 0.4,
         "Meet Gajjar  ·  Aashish Patel  ·  Daksh Gupta",
         size=14, color=INK_SOFT, bold=True, font=FONT_MONO)
    text(s, 0.6, 6.05, 10, 0.4,
         "Three engineers. Three language models. One matching platform.",
         size=13, color=INK_FAINT, italic=True)

    footer(s, 1)
    notes(s, """
OPENING — any presenter, 20 seconds.

"Good morning. We're CureMatch — Meet, Aashish, and Daksh. Our project is
called Matching Clinical Trials With Large Language Models. Over the next
twelve minutes we'll show you the problem, the system we built, and the
three language models we each built inside it."

Pace: warm, confident. Don't rush. Point at the deck.
""")

    # ════════════════════════════════════════════════════════════
    # 02 · THE PROBLEM
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "The Problem")

    title(s, "80% of trials miss\ntheir enrollment deadlines.",
          top=1.2, height=2.2, size=48, line_spacing=1.05, letter_spacing=-30)

    text(s, 0.6, 4.3, 9.5, 1.6,
         "Meanwhile 300,000 patients a year look for a trial and give up.\n"
         "Every trial publishes its eligibility as a wall of medical prose.\n"
         "65,081 actively recruiting studies. No patient reads 65,081 trials.",
         size=17, color=INK_SOFT, line_spacing=1.6)

    rect(s, 0.6, 6.2, 9.5, 0.6, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, 0.85, 6.28, 9.1, 0.22, "A TYPICAL ELIGIBILITY CRITERION",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    text(s, 0.85, 6.5, 9.1, 0.25,
         "ECOG 0–1 · eGFR ≥ 30 · no prior anti-PD-1 · no active autoimmune disease · HbA1c < 10.5%",
         size=10.5, color=INK_SOFT, font=FONT_MONO)

    footer(s, 2)
    notes(s, """
THE PROBLEM — 45 seconds.

Key points to say:
• "80% of clinical trials miss their enrollment deadlines — that's from
   industry research on trial dropout."
• "Meanwhile 300,000 patients a year look for a trial and give up because
   the eligibility is a wall of medical prose."
• "65,081 trials are actively recruiting right now. That's what we built
   CureMatch against. No patient reads 65k trials."
• Point at the bottom example: "This is one criterion from one trial.
   ECOG, eGFR, no prior anti-PD-1, HbA1c — it's written for clinicians."

Transition: "So we built a system to read all 65 thousand for you."
""")

    # ════════════════════════════════════════════════════════════
    # 03 · THE CORPUS — phase donut + top conditions (CHART SLIDE)
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "The Corpus")

    title(s, "65,081 trials. All real.",
          size=40, height=0.9, letter_spacing=-24)

    text(s, 0.6, 2.0, 10, 0.4,
         "Ingested live from ClinicalTrials.gov — the public-domain registry from the U.S. NIH.",
         size=13, color=INK_MUTE)

    image(s, CHARTS_DIR / "phase_donut.png",
          left=0.25, top=2.65, width=5.2)
    image(s, CHARTS_DIR / "top_conditions.png",
          left=5.55, top=2.9, width=4.9)

    text(s, 0.6, 6.55, 10, 0.3,
         "Left: phase distribution across the corpus  ·  Right: the ten most-studied conditions by MeSH tag.",
         size=10, color=INK_FAINT, italic=True, align="center")

    footer(s, 3)
    notes(s, """
THE CORPUS — 50 seconds.

Key points:
• "Everything you're about to see runs on real data — 65,081 trials
   ingested from ClinicalTrials.gov, the U.S. public trial registry."
• Donut: "Two-thirds of trials don't declare a phase — that's Phase Zero
   research and observational studies. Phase 2 is the largest declared
   phase at ten percent."
• Bar chart: "Most-studied conditions — Breast Cancer tops the list at
   1,150 trials, followed by Stroke, Obesity, Prostate Cancer. Oncology
   dominates the actively-recruiting corpus."
• "These numbers are live — the dashboard in our app queries them from
   SQLite on every page load."

Transition to slide 4: "Here's what we built on top of this."
""")

    # ════════════════════════════════════════════════════════════
    # 04 · WHAT WE BUILT
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "What We Built")

    title(s, "Read every trial, for you.", size=48, height=1.2)

    text(s, 0.6, 2.6, 10, 0.5,
         "One profile. Every trial screened. Ranked results. Plain-English explanations.",
         size=16, color=INK_MUTE)

    stages = [
        ("01", "Profile",
         "Age · conditions · medications · labs · location. Three minutes. Browser-only."),
        ("02", "Match",
         "Scored against every actively-recruiting trial. Rule-based, deterministic, reproducible."),
        ("03", "Explain",
         "Every criterion shown. Every exclusion named. Narrative from an LLM."),
    ]
    w, gap, top, h = 3.1, 0.15, 4.2, 2.4
    for i, (n, t, b) in enumerate(stages):
        x = 0.6 + i * (w + gap)
        rect(s, x, top, w, h, WHITE, line=LINE_SOFT, rounded=True)
        text(s, x + 0.3, top + 0.3, w - 0.6, 0.6, n,
             size=36, color=ACCENT, bold=True, font=FONT_MONO, letter_spacing=-20)
        text(s, x + 0.3, top + 1.0, w - 0.6, 0.5, t,
             size=18, color=INK, bold=True)
        text(s, x + 0.3, top + 1.5, w - 0.6, 0.9, b,
             size=11.5, color=INK_MUTE, line_spacing=1.45)

    footer(s, 4)
    notes(s, """
WHAT WE BUILT — 40 seconds.

Read top-down, left to right:
• "CureMatch is three steps for the patient."
• "Step one — enter your profile. Age, conditions, medications, lab
   values, location. Takes three minutes. Stays in your browser; nothing
   uploaded."
• "Step two — match. Every profile is scored against every trial using
   a rule-based engine. Deterministic — same profile always returns the
   same ranked trials."
• "Step three — explain. For every match, we show the full criterion
   breakdown and an LLM writes a plain-English narrative."

Transition: "Here's the architecture that makes this tick."
""")

    # ════════════════════════════════════════════════════════════
    # 05 · ARCHITECTURE
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "The Architecture")

    title(s, "LLMs read. Rules decide. Humans verify.",
          size=34, height=1.0, letter_spacing=-22)

    text(s, 0.6, 2.15, 10, 0.5,
         "Three layers, each doing the one thing it's best at.",
         size=15, color=INK_MUTE)

    cols = [
        ("LLMs READ",
         "Eligibility prose → structured fields.\nPatient questions → grounded answers.\nVerdicts → plain-English narrative.",
         ACCENT),
        ("RULES DECIDE",
         "Age · gender · condition · medications · labs · proximity\n→ weighted composite score.\nSame input, same output, every time.",
         INK),
        ("HUMANS VERIFY",
         "Trial investigators make the final call.\nOur output is ranked candidates — never a diagnosis.\nEvery criterion shown by name.",
         WARNING),
    ]
    col_w = 3.1
    col_gap = 0.15
    col_top = 3.0
    col_h = 3.6
    for i, (label, body, bar_color) in enumerate(cols):
        x = 0.6 + i * (col_w + col_gap)
        rect(s, x, col_top, col_w, col_h, WHITE, line=LINE_SOFT, rounded=True)
        rect(s, x, col_top, col_w, 0.12, bar_color)
        text(s, x + 0.3, col_top + 0.35, col_w - 0.6, 0.4, label,
             size=11, color=bar_color, bold=True, letter_spacing=250)
        text(s, x + 0.3, col_top + 0.95, col_w - 0.6, 2.5, body,
             size=13, color=INK_SOFT, line_spacing=1.6)

    footer(s, 5)
    notes(s, """
ARCHITECTURE — 45 seconds. This is THE core design idea.

Say it exactly:
• "Our entire design philosophy fits on one slide. LLMs read. Rules
   decide. Humans verify."
• Blue column: "LLMs do three language jobs — reading prose into
   structured fields, answering patient questions, narrating verdicts."
• Black column: "Rules do one job — scoring. Deterministic, auditable,
   reproducible. Same input always gives the same output."
• Amber column: "A human always verifies. Our output is ranked
   candidates — it's never a diagnosis. Trial investigators make the
   final call."
• Close the slide: "Medical matching has to be auditable. You can't
   defend an exclusion by saying 'the model felt so.'"

Transition: "Here's the actual system diagram."
""")

    # ════════════════════════════════════════════════════════════
    # 06 · ARCHITECTURE DIAGRAM (visual)
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Architecture  ·  at a glance")

    title(s, "Three containers. One request path.",
          size=30, height=0.9, letter_spacing=-20)

    text(s, 0.6, 1.9, 10, 0.4,
         "Browser calls the server, server reads SQLite, server streams LLM responses back.",
         size=12, color=INK_MUTE, italic=True)

    image(s, CHARTS_DIR / "architecture.png",
          left=0.35, top=2.45, width=10.0)

    footer(s, 6)
    notes(s, """
ARCHITECTURE DIAGRAM — 45 seconds.

Walk through the numbered flow:
• "Three primary containers — browser, Next.js server, data layer."
• "Arrow 1 — patient profile leaves the browser as JSON to the server."
• "Arrow 2 — server asks SQLite for candidate trials using SQL."
• "Arrow 3 — rows come back. Scoring engine runs on them."
• "Arrow 4 — ranked matches return to the browser as JSON."
• "Arrow 5 is the dashed purple arrow — only two routes call the LLM:
   match/explain and chat. The rest of the system never touches Groq."
• External services — OpenStreetMap for geocoding, Groq for LLM —
   are dashed because they're cross-network, not local.
• Bottom band — build-time pipelines. These run once when we ingest
   and parse the corpus. Never on the request path.

Transition: "So here's who built each piece."
""")

    # ════════════════════════════════════════════════════════════
    # 07 · TEAM + SPLIT
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "The Team")

    title(s, "Three engineers. Three LLMs.",
          size=38, height=1.0, letter_spacing=-24)

    text(s, 0.6, 2.15, 10, 0.5,
         "Each of us built a distinct language model into the product.",
         size=15, color=INK_MUTE)

    cards = [
        ("Aashish Patel", "LLM #1", "The Match Explainer",
         "Reads deterministic verdicts; writes a 2-sentence patient-facing narrative. Refusal-tuned."),
        ("Meet Gajjar", "LLM #2 + Frontend", "Chat Agent · LoRA fine-tune",
         "RAG chat grounded in each trial's text. Fine-tuned Llama 3.2 1B. The entire Next.js app."),
        ("Daksh Gupta", "LLM #3", "The Eligibility Parser",
         "Reads compound criteria, negation, paraphrase that regex can't. Feeds the same scorer."),
    ]
    w, gap, top, h = 3.1, 0.15, 3.1, 3.65
    for i, (who, tag, llm_title, body) in enumerate(cards):
        x = 0.6 + i * (w + gap)
        rect(s, x, top, w, h, WHITE, line=LINE_SOFT, rounded=True)
        text(s, x + 0.3, top + 0.3, w - 0.6, 0.35, tag,
             size=9.5, color=ACCENT, bold=True, letter_spacing=250)
        text(s, x + 0.3, top + 0.7, w - 0.6, 0.55, who,
             size=19, color=INK, bold=True, letter_spacing=-10)
        text(s, x + 0.3, top + 1.32, w - 0.6, 0.8, llm_title,
             size=13, color=INK_SOFT, bold=True, italic=True, line_spacing=1.35)
        text(s, x + 0.3, top + 2.15, w - 0.6, 1.3, body,
             size=11.5, color=INK_MUTE, line_spacing=1.55)

    footer(s, 7)
    notes(s, """
THE TEAM — 30 seconds.

Say:
• "The project title says language models — plural. There are three of
   them. Each of us built one."
• "Aashish built the Match Explainer — the LLM that writes patient-
   facing narratives."
• "I, Meet, built the Chat Agent and the LoRA fine-tune, plus the whole
   Next.js frontend — the 3D scenes, the globe, the dashboard."
• "Daksh built the LLM Eligibility Parser — the one that reads what
   regex can't."

Transition: "Aashish — take it away."
""")

    # ════════════════════════════════════════════════════════════
    # 07 · AASHISH — MATCH EXPLAINER
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Aashish Patel  ·  LLM #1")

    title(s, "The Match Explainer.", size=44, height=1.1, letter_spacing=-26)

    text(s, 0.6, 2.1, 10, 0.5,
         "A patient clicks \"Explain this match in plain English.\" The LLM narrates — grounded in the verdicts.",
         size=15, color=INK_MUTE, italic=True)

    stages = [
        ("Verdicts",  "deterministic\nscoring output"),
        ("Serialize", "compact criteria\nstring"),
        ("LLM",       "Groq Llama 3.3 70B\nstreaming"),
        ("Narrative", "≤ 80 words\npatient-facing"),
    ]
    sw, sgap, stop, sh = 2.25, 0.15, 3.2, 1.2
    total_w = len(stages) * sw + (len(stages) - 1) * sgap
    sleft = (SLIDE_W_IN - total_w) / 2
    for i, (lbl, sub) in enumerate(stages):
        x = sleft + i * (sw + sgap)
        rect(s, x, stop, sw, sh, WHITE, line=LINE_SOFT, rounded=True)
        text(s, x, stop + 0.25, sw, 0.4, lbl,
             size=15, color=INK, bold=True, align="center")
        text(s, x, stop + 0.65, sw, 0.5, sub,
             size=10.5, color=INK_MUTE, align="center", line_spacing=1.4)
        if i < len(stages) - 1:
            text(s, x + sw + 0.01, stop + 0.45, sgap - 0.02, 0.4, "›",
                 size=20, color=INK_FAINT, align="center")

    rect(s, 0.6, 4.9, 9.5, 1.8, WHITE, line=LINE_SOFT, rounded=True)
    text(s, 0.85, 5.05, 9.0, 0.3, "SAMPLE OUTPUT",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    text(s, 0.85, 5.35, 9.0, 1.3,
         "\u201CYou match on age, condition, and HbA1c — the trial accepts patients 40–75 with "
         "type 2 diabetes and HbA1c in range, which you meet. Your Warfarin use is on the "
         "trial's exclusion list; the protocol requires no current anticoagulant therapy.\u201D",
         size=13, color=INK_SOFT, italic=True, line_spacing=1.6)

    footer(s, 8)
    notes(s, """
AASHISH — THE MATCH EXPLAINER · 50 seconds.

Cover these beats:
• "The Match Explainer runs every time a patient clicks 'Explain this
   match in plain English.'"
• Walk through the pipeline boxes left to right: "First, the scoring
   engine produces deterministic verdicts per criterion. Second, I
   serialize them into a compact string — each line is 'criterion,
   reader value, required value, verdict'. Third, that goes to Groq
   running Llama 3.3 70B. Fourth, the model streams back a two- to
   three-sentence narrative."
• Point at the sample output: "This is the actual output for a diabetes
   patient with Warfarin. It leads with what matches — age, condition,
   HbA1c — then names the exclusion specifically — 'your Warfarin use is
   on the exclusion list.'"
• Key point: "The LLM never changes the match. It only describes it.
   That's how we keep determinism."

Transition: "Here's exactly how I constrained it to do that."
""")

    # ════════════════════════════════════════════════════════════
    # 08 · AASHISH — THE PROMPT
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Aashish Patel  ·  Prompt Engineering")

    title(s, "80 words. Strongest positive first.",
          size=32, height=0.9, letter_spacing=-20)

    text(s, 0.6, 2.0, 10, 0.5,
         "A well-engineered prompt is a contract with the model — tight scope, explicit shape, no room to drift.",
         size=14, color=INK_MUTE, italic=True)

    code_left = 1.55
    code_w = 7.55
    rect(s, code_left, 2.8, code_w, 3.6, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, code_left + 0.3, 2.95, code_w - 0.6, 0.3, "SYSTEM PROMPT (abridged)",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    prompt_lines = (
        "You are the CureMatch explainer.\n"
        "Given per-criterion verdicts (match / excluded /\n"
        "unknown), write ONE paragraph:\n\n"
        "  • 2–3 sentences, ≤ 80 words.\n"
        "  • Lead with the strongest positive match.\n"
        "  • Name exclusions specifically — e.g. \"Warfarin\n"
        "    is on the exclusion list\" — never generically.\n"
        "  • Close factually. No hype, no uncertainty\n"
        "    beyond what the verdicts show.\n"
        "  • Never fabricate criteria."
    )
    text(s, code_left + 0.3, 3.25, code_w - 0.6, 3.1, prompt_lines,
         size=11, color=INK, font=FONT_MONO, line_spacing=1.5)

    text(s, 0.6, 6.55, 10, 0.3,
         "Result: deterministic verdicts stay deterministic. The LLM never changes the match — it only translates.",
         size=11, color=INK_FAINT, italic=True, align="center")

    footer(s, 9)
    notes(s, """
AASHISH — PROMPT ENGINEERING · 45 seconds.

Walk through the prompt like a contract:
• "A prompt is a contract with the model — it says 'here is exactly
   what I expect.'"
• "First rule — word cap. Two to three sentences, under 80 words.
   Keeps it skimmable."
• "Second — lead with the strongest positive. Patient reads the first
   line first. Start with why they might match."
• "Third — name exclusions specifically. Not 'you have an exclusion.'
   'Your Warfarin use is on the exclusion list.'"
• "Fourth — close factually. No marketing language, no false certainty."
• "Last line, the safety clause — never fabricate criteria. If it's not
   in the verdicts, it doesn't go in the narrative."
• Finisher: "That's the whole prompt. Fifteen lines. But those
   fifteen lines are what make the LLM's output trustworthy."

Transition: "Meet — the frontend and the rest of the LLM layer."
""")

    # ════════════════════════════════════════════════════════════
    # 09 · MEET — THE APP
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Meet Gajjar  ·  Frontend")

    title(s, "The whole application.", size=44, height=1.1, letter_spacing=-26)

    text(s, 0.6, 2.1, 10, 0.5,
         "Next.js 14 · React Three Fiber · GSAP scroll storytelling · Leaflet maps · Recharts dashboard.",
         size=14, color=INK_MUTE)

    features = [
        ("Scroll storytelling", "/   landing page",
         "DNA helix, library, scan, verdict — GSAP ScrollTrigger narrates the product before a button is clicked."),
        ("3D results globe", "/results",
         "All matched trial sites plotted on a pulsing R3F globe. Geocoded via Nominatim + haversine proximity."),
        ("Live corpus dashboard", "/data",
         "Recharts visualizing 65,081 trials: phase distribution, top conditions, parser coverage, excluded-med frequencies."),
        ("Trial detail + chat", "/trial/[id]",
         "Match breakdown with real patient values vs trial requirements. Floating chat launcher. Streaming narratives."),
    ]
    w, h, gap = 4.7, 1.25, 0.15
    base_top = 3.0
    for i, (t, route, body) in enumerate(features):
        row, col = i // 2, i % 2
        x = 0.6 + col * (w + gap)
        y = base_top + row * (h + gap)
        rect(s, x, y, w, h, WHITE, line=LINE_SOFT, rounded=True)
        text(s, x + 0.25, y + 0.2, w - 2.0, 0.4, t,
             size=14, color=INK, bold=True)
        text(s, x + w - 2.0, y + 0.22, 1.8, 0.35, route,
             size=10, color=ACCENT, font=FONT_MONO, align="right")
        text(s, x + 0.25, y + 0.62, w - 0.45, 0.6, body,
             size=10.5, color=INK_MUTE, line_spacing=1.45)

    footer(s, 10)
    notes(s, """
MEET — THE APP · 50 seconds.

Say what you built:
• "Before the LLMs, I'll briefly show you what I built as the whole
   application. Next.js 14 on the frontend, SQLite on the backend."
• Four features top to bottom, left to right:
   1. "Scroll storytelling — the landing page is a GSAP scroll scene
      with 3D DNA helix, a library, a scanning animation, a verdict.
      The product explains itself before you click anything."
   2. "3D results globe — every matched trial site is plotted on a
      pulsing Three.js globe. Geocoded with Nominatim and haversine
      proximity to rank closest sites."
   3. "Live corpus dashboard — at /data, the page you saw on slide 3.
      Recharts over the real SQLite numbers."
   4. "Trial detail page — per-criterion breakdown with real patient
      values against real requirements. A floating chat launcher at the
      bottom right. Both narratives stream in real-time."
• "All of this stays in the browser. No accounts, no tracking."

Transition: "Now the first LLM I built — the chat agent."
""")

    # ════════════════════════════════════════════════════════════
    # 10 · MEET — CHAT AGENT
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Meet Gajjar  ·  LLM #2  ·  The Chat Agent")

    title(s, "Ask. Grounded. Refused when off-topic.",
          size=30, height=0.9, letter_spacing=-20)

    text(s, 0.6, 2.0, 10, 0.4,
         "Llama 3.3 70B via Groq · RAG over each trial's eligibility text · streaming at ~500 tokens/sec.",
         size=13, color=INK_MUTE)

    ex_left, ex_w, ex_top, ex_h = 0.6, 4.7, 2.75, 4.15

    rect(s, ex_left, ex_top, ex_w, ex_h, WHITE, line=LINE_SOFT, rounded=True)
    text(s, ex_left + 0.25, ex_top + 0.18, ex_w - 0.5, 0.3, "IN-CONTEXT",
         size=9, color=SUCCESS, bold=True, letter_spacing=200)
    text(s, ex_left + 0.25, ex_top + 0.5, ex_w - 0.5, 0.4,
         "\u201CCan I join if I'm on Warfarin?\u201D",
         size=13, color=INK, italic=True, bold=True)
    text(s, ex_left + 0.25, ex_top + 1.02, ex_w - 0.5, 1.3,
         "The trial's eligibility criteria list anticoagulants in the exclusions, so you "
         "would not qualify while taking Warfarin.",
         size=11.5, color=INK_SOFT, line_spacing=1.55)

    rect(s, ex_left + 0.25, ex_top + 2.3, ex_w - 0.5, 0.02, LINE_SOFT)

    text(s, ex_left + 0.25, ex_top + 2.5, ex_w - 0.5, 0.3, "OUT-OF-CONTEXT",
         size=9, color=WARNING, bold=True, letter_spacing=200)
    text(s, ex_left + 0.25, ex_top + 2.85, ex_w - 0.5, 0.4,
         "\u201CWhat's the weather in Tokyo?\u201D",
         size=13, color=INK, italic=True, bold=True)
    text(s, ex_left + 0.25, ex_top + 3.35, ex_w - 0.5, 0.7,
         "The trial's public information doesn't specify that.",
         size=11.5, color=INK_SOFT, line_spacing=1.55)

    how_left, how_w = 5.55, 4.55
    rect(s, how_left, ex_top, how_w, ex_h, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, how_left + 0.3, ex_top + 0.25, how_w - 0.6, 0.3, "HOW IT WORKS",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)

    steps = [
        ("1",  "Load trial's eligibility text from trials.db."),
        ("2",  "Inject it as retrieval context in the prompt."),
        ("3",  "Forbid answering outside the text; script refusal."),
        ("4",  "Append the user turn; stream from Groq."),
        ("5",  "Persist history per-NCT in localStorage."),
    ]
    step_gap = 0.66
    for i, (n, t) in enumerate(steps):
        y = ex_top + 0.8 + i * step_gap
        text(s, how_left + 0.3, y, 0.5, 0.4, n,
             size=15, color=ACCENT, bold=True, font=FONT_MONO)
        text(s, how_left + 0.75, y + 0.02, how_w - 0.95, 0.55, t,
             size=12, color=INK_SOFT, line_spacing=1.4)

    footer(s, 11)
    notes(s, """
MEET — THE CHAT AGENT · 55 seconds.

Show both exchanges, then explain the architecture.

Left card — the behavior:
• "Here's an in-context question. Patient asks 'Can I join if I'm on
   Warfarin?' — the agent answers from the trial's actual eligibility
   text, citing the exclusion directly."
• "And here's an out-of-context question — 'What's the weather in
   Tokyo?' — the agent refuses. Not 'I don't know.' Specifically
   'The trial's public information doesn't specify that.'"

Right card — how it works:
• "Five steps. One — load the trial's eligibility text. Two — inject
   it as retrieval context. Three — the system prompt forbids
   answering outside that text and scripts the refusal. Four — append
   the user turn and stream from Groq. Five — persist history per-NCT
   in localStorage so the conversation resumes next visit."
• "The refusal is load-bearing. Without it a chat agent on medical
   data is a liability. With it we have a product."

Transition: "I also fine-tuned a smaller model to verify this works."
""")

    # ════════════════════════════════════════════════════════════
    # 11 · MEET — LoRA FINE-TUNE
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Meet Gajjar  ·  Fine-Tune")

    title(s, "LoRA on Llama 3.2 1B.",
          size=40, height=1.0, letter_spacing=-24)

    text(s, 0.6, 2.05, 10, 0.4,
         "A controlled experiment: parser-grounded Q&A improves refusal behavior in a small model.",
         size=13, color=INK_MUTE)

    stats = [
        ("450",     "Q&A pairs auto-generated from parser output"),
        ("rank 16", "LoRA · 8M adapter params on 1B base"),
        ("3 epochs","Colab T4 GPU · ~40 min runtime"),
        ("$0",      "total training cost · free tier"),
    ]
    card_left, card_w, card_h, card_top = 0.6, 4.55, 0.82, 2.75
    for i, (n, label) in enumerate(stats):
        y = card_top + i * (card_h + 0.12)
        rect(s, card_left, y, card_w, card_h, WHITE, line=LINE_SOFT, rounded=True)
        text(s, card_left + 0.3, y + 0.14, 1.8, 0.6, n,
             size=22, color=ACCENT, bold=True, font=FONT_MONO, letter_spacing=-15)
        text(s, card_left + 2.15, y + 0.26, card_w - 2.3, 0.5, label,
             size=12, color=INK_MUTE, line_spacing=1.45)

    loss = PROJECT_ROOT / "loss_curve.png"
    if not loss.exists():
        loss = CHARTS_DIR / "loss_curve.png"
    if loss.exists():
        image(s, loss, left=5.45, top=2.75, width=4.7)
        text(s, 5.45, 6.4, 4.7, 0.3,
             "Training + eval loss across 3 epochs  ·  notebooks/finetune_curematch.ipynb",
             size=9, color=INK_FAINT, italic=True, align="center")

    text(s, 0.6, 6.85, 10, 0.3,
         "Production inference uses base Llama 3.3 70B via Groq. The 1B fine-tune proves the dataset-construction pipeline works.",
         size=10, color=INK_FAINT, italic=True, align="center")

    footer(s, 12)
    notes(s, """
MEET — FINE-TUNE · 55 seconds.

Say:
• "This is the model-level work. I took Meta's Llama 3.2 1B Instruct
   and fine-tuned it with LoRA — Low-Rank Adaptation."
• Walk through the stats: "450 Q&A pairs. I auto-generated these from
   our rule-parser output — every example is grounded in a real trial.
   Rank 16 LoRA — 8 million trainable parameters on a 1 billion
   parameter base. 3 epochs on a free Colab T4. 40 minutes total. Zero
   dollars."
• Point at loss curve: "This is the training curve. Blue is training
   loss, orange is eval loss. Both converge cleanly across three epochs.
   Final train loss 0.43, eval 0.70 — no overfitting."
• The honest framing: "In production we use base Llama 3.3 70B through
   Groq, because it's faster and we didn't need to ship the 1B. The
   fine-tune was a controlled experiment — it proves our dataset-
   construction pipeline works. If we scaled up, that pipeline scales."

Transition: "Daksh — the LLM that reads the trials."
""")

    # ════════════════════════════════════════════════════════════
    # 12 · DAKSH — LLM PARSER
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Daksh Gupta  ·  LLM #3  ·  The Parser")

    title(s, "Reading what regex can't.",
          size=40, height=1.0, letter_spacing=-24)

    text(s, 0.6, 2.05, 10, 0.5,
         "The rule parser covers 33% of trials. For the other 67%, the LLM reads compound criteria, negation, and paraphrase.",
         size=13, color=INK_MUTE, line_spacing=1.5)

    rect(s, 0.6, 3.0, 9.5, 0.9, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, 0.85, 3.12, 9.0, 0.25, "INPUT  ·  one line from an eligibility criterion",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    text(s, 0.85, 3.4, 9.0, 0.45,
         "\"Creatinine clearance > 60 mL/min, unless on chronic dialysis, in which case any value is acceptable.\"",
         size=11.5, color=INK_SOFT, font=FONT_MONO)

    col_top, col_h, col_w = 4.1, 2.5, 4.65

    rect(s, 0.6, col_top, col_w, col_h, WHITE, line=LINE_SOFT, rounded=True)
    rect(s, 0.6, col_top, col_w, 0.12, RGBColor(0xFF, 0x3B, 0x30))
    text(s, 0.85, col_top + 0.3, col_w - 0.5, 0.3, "RULE PARSER",
         size=10, color=RGBColor(0xFF, 0x3B, 0x30), bold=True, letter_spacing=200)
    text(s, 0.85, col_top + 0.72, col_w - 0.5, 0.5, "{ missed the override }",
         size=14, color=INK, bold=True, font=FONT_MONO)
    text(s, 0.85, col_top + 1.25, col_w - 0.5, 1.1,
         "Matched \"creatinine > 60\" only. Missed the dialysis conditional — dialysis "
         "patients would be incorrectly excluded from the trial.",
         size=11, color=INK_MUTE, line_spacing=1.5)

    rect(s, 5.45, col_top, col_w, col_h, WHITE, line=LINE_SOFT, rounded=True)
    rect(s, 5.45, col_top, col_w, 0.12, SUCCESS)
    text(s, 5.7, col_top + 0.3, col_w - 0.5, 0.3, "LLM PARSER",
         size=10, color=SUCCESS, bold=True, letter_spacing=200)
    text(s, 5.7, col_top + 0.72, col_w - 0.5, 0.5, "{ extracted both }",
         size=14, color=INK, bold=True, font=FONT_MONO)
    text(s, 5.7, col_top + 1.25, col_w - 0.5, 1.1,
         "lab: creatinine_clearance, min: 60\n"
         "override: { if: on_dialysis, ignore_threshold: true }",
         size=11, color=INK_SOFT, line_spacing=1.6, font=FONT_MONO)

    footer(s, 13)
    notes(s, """
DAKSH — THE LLM PARSER · 50 seconds.

Say:
• "I built the third LLM — the eligibility parser. It reads what
   regex can't."
• Point at the input: "Here's a real line from a real trial's
   eligibility: 'Creatinine clearance greater than 60, unless on chronic
   dialysis, in which case any value is acceptable.' That's a compound
   criterion with a conditional override."
• Left card: "The rule parser saw 'creatinine greater than 60' and
   stopped there. It missed the override. Dialysis patients would be
   incorrectly excluded."
• Right card: "The LLM reads the full sentence. It extracts the
   threshold — creatinine minimum 60 — AND the override rule. Same
   structure the scoring engine expects downstream."
• Key point: "The output of the LLM is fed into the SAME deterministic
   scorer. We don't lose auditability — we gain coverage."

Transition: "Here's how much coverage we gain."
""")

    # ════════════════════════════════════════════════════════════
    # 13 · DAKSH — PARSER COVERAGE (CHART SLIDE)
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Daksh Gupta  ·  Parser Coverage")

    title(s, "The 33% ceiling.",
          size=36, height=0.9, letter_spacing=-22)

    text(s, 0.6, 1.95, 10, 0.4,
         "Rules hit a wall. Compound criteria, negation, paraphrase need language understanding.",
         size=13, color=INK_MUTE)

    image(s, CHARTS_DIR / "parser_coverage.png",
          left=1.3, top=2.5, height=3.5)

    rect(s, 0.6, 6.2, 9.5, 0.75, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, 0.85, 6.32, 9.0, 0.28, "THE OPPORTUNITY",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    text(s, 0.85, 6.6, 9.0, 0.3,
         "The 67% gap is what an LLM parser closes. Same scorer downstream — triple the coverage.",
         size=12, color=ACCENT, bold=True)

    footer(s, 14)
    notes(s, """
DAKSH — PARSER COVERAGE · 40 seconds.

Chart is self-explanatory, narrate the numbers:
• "Here's the coverage the rule-parser achieves across the full 65,081
   trial corpus. Medication exclusions — 32.6 percent. Lab thresholds —
   29.3 percent. ECOG performance status — 4 percent."
• "86 seconds to parse all 65k trials. Zero cost to operate."
• The KEY takeaway: "But 67 percent of trials slip through. That's
   the ceiling — the limit of what regex and drug dictionaries can
   reach."
• Point at the blue footer line: "That 67 percent gap is the opportunity.
   Swap the rule parser for an LLM parser, feed the same scoring engine,
   triple the coverage. Schema's already ready for it — the parsed
   database has a source column for 'rules' versus 'llm'."

Transition: "Because the scorer doesn't care where the fields came from."
""")

    # ════════════════════════════════════════════════════════════
    # 14 · DAKSH — SCORING ENGINE
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "Daksh Gupta  ·  Scoring Engine")

    title(s, "LLM reads. Rule engine scores.",
          size=34, height=1.0, letter_spacing=-22)

    text(s, 0.6, 2.05, 10, 0.5,
         "Whether the structured fields came from rules or an LLM, the downstream scorer is the same — deterministic, auditable.",
         size=13, color=INK_MUTE, line_spacing=1.5)

    weights = [
        ("Condition",    "35%"),
        ("Medications",  "20%"),
        ("Lab values",   "15%"),
        ("Age",          "12%"),
        ("Proximity",    "10%"),
        ("Gender",        "8%"),
    ]
    tbl_top = 3.0
    row_h = 0.45
    for i, (name, pct) in enumerate(weights):
        y = tbl_top + i * row_h
        text(s, 1.5, y, 3.5, 0.4, name, size=16, color=INK_SOFT)
        text(s, 5.2, y, 2.0, 0.4, pct,
             size=16, color=ACCENT, bold=True, font=FONT_MONO)
        rect(s, 1.5, y + 0.4, 8.2, 0.01, LINE_SOFT)

    rect(s, 0.6, 5.95, 9.5, 0.85, PAPER_ALT, line=LINE_SOFT, rounded=True)
    text(s, 0.85, 6.08, 9.0, 0.25, "AUDITABLE EXCLUSION",
         size=9, color=INK_FAINT, bold=True, letter_spacing=200)
    text(s, 0.85, 6.34, 9.0, 0.45,
         "You (HbA1c 6.2)  →  Required (≥ 7.5)  →  excluded.",
         size=16, color=INK, bold=True, font=FONT_MONO, letter_spacing=-8)

    footer(s, 15)
    notes(s, """
DAKSH — SCORING ENGINE · 40 seconds.

Walk through the weights:
• "The scoring engine runs six rule-based checks per trial. Condition
   match is 35 percent — we weight that highest because condition
   mismatch is the worst kind of false positive."
• "Medications 20, labs 15, age 12, proximity 10, gender 8."
• "Composite is a weighted average. Any trial excluded on condition is
   dropped entirely — we don't surface a prostate cancer trial to a
   breast cancer patient, no matter what."
• Point at the audit line: "And this is what auditable means. If a
   hospital asks 'why was this patient excluded?' we point at a rule.
   'You had HbA1c 6.2. Trial required at least 7.5. Excluded.' Same
   answer every time."
• Closer: "The LLM helps us read more trials. The rules make sure the
   matching stays defensible."

Transition: "From 65k trials to the top five — in milliseconds."
""")

    # ════════════════════════════════════════════════════════════
    # 15 · END-TO-END FUNNEL (CHART SLIDE)
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "The Matching Funnel")

    title(s, "65,081 → 5. In milliseconds.",
          size=36, height=0.9, letter_spacing=-22)

    text(s, 0.6, 1.95, 10, 0.4,
         "Every profile goes through the same four-stage filter. The scoring only touches ~2% of the corpus.",
         size=13, color=INK_MUTE)

    image(s, CHARTS_DIR / "funnel.png",
          left=2.3, top=2.5, width=6.0)

    footer(s, 16)
    notes(s, """
FUNNEL · 35 seconds.

Narrate the funnel:
• "Here's what happens end-to-end when a patient submits their profile."
• "We start with 65,081 trials. SQL pre-filter by distinctive condition
   words narrows that to around 1,572 candidates — for this example, a
   diabetes patient. That step takes a handful of milliseconds in
   SQLite."
• "The scoring engine runs the six-criterion check across those 1,572.
   That produces 147 scored candidates — after we drop hard exclusions."
• "We return the top five to the UI, ranked by composite score."
• "Total wall-clock time: around 250 milliseconds on a laptop. That's
   the whole funnel."

Transition: "So what did we ship?"
""")

    # ════════════════════════════════════════════════════════════
    # 16 · IN NUMBERS
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "In Numbers")

    title(s, "What we shipped.",
          size=44, height=1.1, letter_spacing=-26)

    stats_grid = [
        ("65,081",  "trials indexed"),
        ("411,042", "sites with lat/lng"),
        ("21,191",  "parsed exclusion lists"),
        ("19,062",  "parsed lab thresholds"),
        ("86 sec",  "full-corpus parse time"),
        ("$0",      "ongoing cost"),
    ]
    w, h, gap, top = 3.1, 1.35, 0.15, 3.0
    for i, (n, label) in enumerate(stats_grid):
        row, col = i // 3, i % 3
        x = 0.6 + col * (w + gap)
        y = top + row * (h + 0.18)
        rect(s, x, y, w, h, WHITE, line=LINE_SOFT, rounded=True)
        text(s, x + 0.25, y + 0.22, w - 0.5, 0.7, n,
             size=30, color=ACCENT, bold=True, letter_spacing=-18, font=FONT_MONO)
        text(s, x + 0.25, y + 0.9, w - 0.5, 0.4, label,
             size=11, color=INK_MUTE, line_spacing=1.3)

    text(s, 0.6, 6.1, 10, 0.4,
         "Zero accounts · zero tracking · zero analytics · zero telemetry.",
         size=13, color=INK_MUTE, align="center", italic=True)

    footer(s, 17)
    notes(s, """
IN NUMBERS · 30 seconds.

Read left-to-right, top-down:
• "65,081 trials indexed. 411 thousand trial sites with latitude and
   longitude. 21 thousand parsed exclusion lists. 19 thousand parsed
   lab thresholds. 86 seconds to parse the entire corpus. Zero dollars
   ongoing cost."
• Closer: "And zero accounts, zero tracking, zero analytics, zero
   telemetry. The whole app is free for patients to use."

Transition: "We want to be honest about what we haven't done."
""")

    # ════════════════════════════════════════════════════════════
    # 17 · LIMITATIONS
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, PAPER)
    eyebrow(s, "What's Not Solved", color=WARNING)

    title(s, "Limitations.",
          size=44, height=1.1, letter_spacing=-26)

    limits = [
        ("Research prototype.",
         "Not a medical device. Final eligibility is the trial investigators' call."),
        ("Parser coverage is partial.",
         "33% of trials for medications, 29% for labs. The LLM parser is the planned remediation."),
        ("LLM parser not yet in production.",
         "Schema and source column are ready. We've run it on samples, not the full corpus."),
        ("No published accuracy metric.",
         "Qualitative spot-checks only. An F1 score against hand-labeled ground truth is the next milestone."),
    ]
    item_top = 2.6
    for i, (t_, body) in enumerate(limits):
        y = item_top + i * 0.95
        rect(s, 0.6, y + 0.1, 0.14, 0.14, WARNING, rounded=True, radius=0.5)
        text(s, 0.95, y - 0.05, 9, 0.4, t_,
             size=16, color=INK, bold=True)
        text(s, 0.95, y + 0.35, 9, 0.55, body,
             size=12.5, color=INK_MUTE, line_spacing=1.5)

    footer(s, 18)
    notes(s, """
LIMITATIONS · 35 seconds.

Be honest — it builds credibility:
• "CureMatch is a research prototype. Not a medical device. Every
   narrative ends with 'final eligibility is determined by the trial's
   investigators.' That's enforced in the prompt."
• "Parser coverage is partial — 33 percent for medications, 29 percent
   for labs. The LLM parser is the planned fix, we haven't run it on
   the full corpus yet."
• "We've spot-checked the parser qualitatively. No published F1 against
   hand-labeled ground truth yet — that's the next milestone."
• "And to be clear — we haven't deployed to users. This is a working
   prototype, not a launched product."

Transition: "To close."
""")

    # ════════════════════════════════════════════════════════════
    # 18 · CLOSING
    # ════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, DEEP)

    text(s, 0.6, 1.6, 10, 1.1, "LLMs for language.",
         size=54, color=WHITE, bold=True, letter_spacing=-28)
    text(s, 0.6, 2.55, 10, 1.1, "Rules for decisions.",
         size=54, color=ACCENT_LT, bold=True, letter_spacing=-28)
    text(s, 0.6, 3.5, 10, 1.1, "Humans verify.",
         size=54, color=RGBColor(0xBF, 0x5A, 0xF2), bold=True, letter_spacing=-28)

    rect(s, 0.6, 4.9, 0.5, 0.04, ACCENT_LT)
    text(s, 0.6, 5.15, 10, 0.45, "CureMatch.",
         size=22, color=WHITE, bold=True)
    text(s, 0.6, 5.55, 10, 0.4, "github.com/meetgajjarx07/curematch",
         size=13, color=RGBColor(0x86, 0x86, 0x8B), font=FONT_MONO)
    text(s, 0.6, 6.2, 10, 0.5, "Thank you. Questions?",
         size=18, color=RGBColor(0xC6, 0xC6, 0xCB))

    footer(s, 19, dark=True)
    notes(s, """
CLOSING · 20 seconds.

Land the thesis cleanly. All three presenters stand together for Q&A.

Say:
• "Three principles hold the system together. LLMs for language — they
   read prose and answer questions. Rules for decisions — because
   medical matching has to be auditable. Humans verify — trial
   investigators always make the final call."
• "That separation is the whole product."
• "CureMatch is open-source at the URL shown. Thank you — we'd love
   to take your questions."

Q&A RULES:
• Whoever's LLM is being asked about fields the answer. The other two
   can chime in if they want.
• If it's a deep-technical question, Meet takes it.
• If it's a design-choice question, any of us can answer.
• DO NOT over-promise anything. Research prototype.
""")

    # ─── Save ───────────────────────────────────────────────────
    prs.save(out_path)
    print(f"\n✅ Presentation saved to: {out_path}")
    print(f"   {prs.slide_width / 914400:.2f}\" × {prs.slide_height / 914400:.2f}\" (16:9)")
    print(f"   {len(prs.slides)} slides  ·  speaker notes on every slide")


if __name__ == "__main__":
    out = Path(__file__).resolve().parent.parent / "CureMatch_Presentation.pptx"
    build(out)
